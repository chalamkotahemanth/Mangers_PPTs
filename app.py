import io
import re
import zipfile
from datetime import datetime

import pandas as pd
import streamlit as st
from pptx import Presentation

st.set_page_config(page_title="PPT KPI Extractor", page_icon="📊", layout="wide")

# ---------- Regex patterns for KPI extraction ----------
PATTERNS = {
    "Achievement Rate": re.compile(r"Achievement\s*Rate[:\s]*([\d\.]+%)", re.I),
    "Revenue Target": re.compile(r"Revenue\s*Target[:\s]*[₹Rs\.\s]*([\d,]+)", re.I),
    "Achieved": re.compile(r"Achieved[:\s]*[₹Rs\.\s]*([\d,]+)", re.I),
    "Revenue Reached": re.compile(r"revenue\s*reached\s*[₹Rs\.\s]*([\d,]+)", re.I),
    "Target Of": re.compile(r"against\s*a\s*target\s*of\s*[₹Rs\.\s]*([\d,]+)", re.I),
    "Revenue %": re.compile(r"against\s*a\s*target\s*of\s*[₹Rs\.\s]*[^\(]+\(([\d\.]+%)\)", re.I),
    "Quality Score": re.compile(r"Quality\s*score\s*was\s*([\d\.]+%)", re.I),
    # Fallback phrasings
    "Quality Score (alt)": re.compile(r"Quality\s*Score[:\s]*([\d\.]+%)", re.I),
}


def extract_runs_text(shape):
    """Merge text runs so symbols like ₹ + number are captured together."""
    parts = []
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                t = (run.text or "").strip()
                if t:
                    parts.append(t)
    return " ".join(parts).strip()


def parse_percent(val):
    if not val:
        return None
    try:
        return float(val.replace("%", "").strip())
    except Exception:
        return None


def parse_currency(val):
    if not val:
        return None
    try:
        return float(val.replace(",", "").replace(" ", ""))
    except Exception:
        return None


def extract_kpis_from_text(full_text):
    data = {}
    for key, pattern in PATTERNS.items():
        m = pattern.search(full_text)
        if m:
            data[key] = m.group(1)
    # Normalize common fields
    # Prefer "Quality Score" if alt exists
    if "Quality Score" not in data and "Quality Score (alt)" in data:
        data["Quality Score"] = data["Quality Score (alt)"]
    data.pop("Quality Score (alt)", None)
    return data


def process_pptx_bytes(ppt_bytes, filename):
    prs = Presentation(io.BytesIO(ppt_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            txt = extract_runs_text(shape)
            if txt:
                texts.append(txt)
    full_text = " ".join(texts)
    kpis = extract_kpis_from_text(full_text)
    row = {
        "PPT File": filename,
        "Achievement Rate (%)": parse_percent(kpis.get("Achievement Rate")),
        "Revenue Target (₹)": parse_currency(kpis.get("Revenue Target")),
        "Achieved (₹)": parse_currency(kpis.get("Achieved")),
        "Revenue Reached (₹)": parse_currency(kpis.get("Revenue Reached")),
        "Target Of (₹)": parse_currency(kpis.get("Target Of")),
        "Revenue % (%)": parse_percent(kpis.get("Revenue %")),
        "Quality Score (%)": parse_percent(kpis.get("Quality Score")),
    }
    return row


st.title("📊 PPT KPI Extractor")
st.write(
    "Upload multiple **.pptx** files or a **.zip** containing PPTX files. We'll extract KPIs and give you an Excel to share with your manager."
)

col1, col2 = st.columns(2)
with col1:
    uploaded_ppts = st.file_uploader("Upload PPTX files", type=["pptx"], accept_multiple_files=True)
with col2:
    uploaded_zip = st.file_uploader("...or upload a ZIP containing PPTX files", type=["zip"], accept_multiple_files=False)

process_btn = st.button("🚀 Process Files")

if process_btn:
    rows = []
    seen = set()

    # Handle PPTX files
    if uploaded_ppts:
        for f in uploaded_ppts:
            try:
                rows.append(process_pptx_bytes(f.read(), f.name))
                seen.add(f.name)
            except Exception as e:
                st.warning(f"Failed to process {f.name}: {e}")

    # Handle ZIP
    if uploaded_zip is not None:
        try:
            zf = zipfile.ZipFile(uploaded_zip)
            for name in zf.namelist():
                if name.lower().endswith(".pptx") and name not in seen:
                    try:
                        rows.append(process_pptx_bytes(zf.read(name), name))
                        seen.add(name)
                    except Exception as e:
                        st.warning(f"Failed to process {name} from ZIP: {e}")
        except Exception as e:
            st.error(f"Could not read ZIP: {e}")

    if not rows:
        st.info("No valid PPTX files found.")
    else:
        df = pd.DataFrame(rows)

        # Derived columns for visualization
        # "Best Achieved" prefers 'Achieved' else 'Revenue Reached'
        df["Best Achieved (₹)"] = df.apply(
            lambda r: r["Achieved (₹)"] if pd.notnull(r["Achieved (₹)"]) else r["Revenue Reached (₹)"], axis=1
        )
        # "Best Target" prefers 'Revenue Target' else 'Target Of'
        df["Best Target (₹)"] = df.apply(
            lambda r: r["Revenue Target (₹)"] if pd.notnull(r["Revenue Target (₹)"]) else r["Target Of (₹)"], axis=1
        )
        # Compute Achievement Rate if missing
        def infer_rate(row):
            if pd.notnull(row["Achievement Rate (%)"]):
                return row["Achievement Rate (%)"]
            if (
                pd.notnull(row["Best Achieved (₹)"])
                and pd.notnull(row["Best Target (₹)"])
                and row["Best Target (₹)"] != 0
            ):
                return round((row["Best Achieved (₹)"] / row["Best Target (₹)"]) * 100.0, 2)
            if pd.notnull(row["Revenue % (%)"]):
                return row["Revenue % (%)"]
            return None

        df["Achievement Rate (%)"] = df.apply(infer_rate, axis=1)

        st.subheader("📋 Extracted KPI Table")
        st.dataframe(df, use_container_width=True)

        # Downloads
        excel_buf = io.BytesIO()
        with pd.ExcelWriter(excel_buf, engine="openpyxl") as writer:
            df.to_excel(writer, index=False, sheet_name="KPIs")
        excel_buf.seek(0)

        st.download_button(
            "⬇️ Download Excel (KPIs)",
            data=excel_buf,
            file_name=f"managers_kpi_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

        csv_buf = df.to_csv(index=False).encode("utf-8")
        st.download_button(
            "⬇️ Download CSV",
            data=csv_buf,
            file_name=f"managers_kpi_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
        )

        st.subheader("📈 Quick Visuals")
        # Simple charts
        try:
            st.bar_chart(df.set_index("PPT File")[
                ["Best Target (₹)", "Best Achieved (₹)"]
            ])
        except Exception:
            pass
        try:
            st.line_chart(df.set_index("PPT File")[
                ["Achievement Rate (%)"]
            ])
        except Exception:
            pass

        st.success("Done!  downloaded Excel use the visuals above.")
